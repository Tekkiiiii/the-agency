#!/usr/bin/env python3
"""
audit_deck.py — mechanical audit for deck-narrative rules.

Takes a .pptx path, reports what a machine can actually check:
  - slide count
  - per-slide largest font size (title/headline size)
  - per-slide title text, word count, explicit line-break count,
    an ESTIMATED wrapped-line count (heuristic, not measured)
  - per-slide verb-test heuristic result (Vietnamese + English aware)
  - table count vs chart count (deck-wide, and per slide)
  - a title dump in sequence, for the horizontal-logic read-through
  - per-slide process-meta scan (Rule 8, reader-only content) — flags
    producer/reviewer-facing marker strings ("TBD", "placeholder", "cần xác
    nhận", "chưa xác định", etc.) anywhere in slide text (shapes + table
    cells), as REVIEW CANDIDATES only — see the loud caveat in Rule 8 below
  - deck-wide FONT INVENTORY — every text run (shapes + table cells) AND
    every reachable chart-text font (title, axis tick labels, legend, data
    labels) is walked and its declared font name recorded. Two sub-checks
    ride on top of the inventory, both MECHANICAL/HIGH confidence (not
    heuristics — a font name string either is or isn't in a fixed list):
      1. SYSTEM-FONT FLAG — Calibri, Arial, Times New Roman, Helvetica,
         Cambria, Verdana, Tahoma anywhere in the deck is a hard fail. These
         are Office/OS defaults; their presence in a shipped institutional
         deliverable is unconditionally wrong, independent of which brand's
         deck this is — see deck-narrative/SKILL.md's font section and your
         project's brand skill's "never use system fonts" rule (if it states
         one). On by default, no flag needed.
      2. UNSET/THEME-INHERITED FONT COUNT — runs/elements with no explicit
         font name resolve from the presentation theme, which is invisible
         in the file and can render differently machine to machine. Reported
         as a count, not a verdict — inheriting the theme is not necessarily
         wrong, but it's always worth knowing.
    A third, OPTIONAL check activates only with --expect "Family A,Family B":
    any font present in the deck that is neither a system font (already
    caught above) nor in the --expect list is flagged as OUTSIDE THE
    EXPECTED SET. This script does NOT hardcode brand font names — it stays
    brand-agnostic per its architecture-layer role; the expected family list
    comes from whichever brand/visual skill is co-loaded, passed in by the
    caller. Without --expect, the inventory is reported with no
    family-membership judgment beyond the system-font check.
    Does NOT check whether a font is actually installed on the machine doing
    the audit — installation is a property of the rendering environment, not
    the deck, and a deck can legitimately target a machine other than the
    one running this script.

What this script does NOT do (by design — see deck-narrative/SKILL.md):
  - It does not judge whether a title states a "so what" implication.
  - It does not judge whether a slide's body proves its title (vertical logic).
  - It does not judge narrative flow (horizontal logic) — it only prints the
    title sequence so a human/agent can read it.
  - It does not verify pre-wiring, footnote citations, or decision-framework
    fit (Challenge>Options>Decision / Goal>Obstacle>Strategy>Action /
    Insight>Action>Outcome) — those require reading slide body content in
    context, not just shape geometry.
  - It does NOT decide whether a process-meta marker hit is actually a
    violation (Rule 8 TEST 1: addressed to audience vs. producer/reviewer).
    The rule (deck-narrative/SKILL.md Rule 8) is that the discriminator is
    WHO THE LINE IS ADDRESSED TO — the deck's actual audience vs. the
    producer/an internal reviewer — not whether a marker string like
    "decision" or "cần xác nhận" appears. This script cannot read addressee.
    A marker hit on a slide that asks the deck's real decision-maker to
    approve something is CORRECT AS-IS under Test 1, not a bug. Never
    auto-strip or auto-fail on a hit here.
  - It does NOT and CANNOT run Rule 8 TEST 2 (see vs. say): even content
    that IS audience-facing may belong spoken by the presenter rather than
    shown on a slide — supporting rationale, evidence tables, methodology,
    the reasoning behind an ask. Whether the audience needs to SEE a given
    block of content or the presenter can just SAY it requires knowing what
    the audience already needs vs. what they'd only ask for if curious —
    there is no mechanical signal for that in a .pptx file. This is pure
    judgment, on every slide, including ones that pass every other check in
    this script. See deck-narrative/SKILL.md Rule 8 for the worked example
    (three decision slides that should compress to one, with the supporting
    research spoken instead of slotted onto slides 5-7).
  - The verb test, the line-wrap estimate, and the process-meta scan are
    HEURISTICS, not ground truth. They are flagged as such in the output.
    Treat every FAIL/REVIEW flag as a candidate for human or agent judgment,
    not a final verdict.
  - Font scan blind spot: chart text is reached at the axis/legend/series
    level (chart.category_axis.tick_labels.font, value_axis.tick_labels.font,
    chart.legend.font, plot.data_labels.font, series.data_labels.font,
    chart_title runs) — confirmed reachable via python-pptx against this
    deck's own charts. What is NOT reached: a font override set on an
    INDIVIDUAL data point's label (point.data_label.font), which python-pptx
    exposes one point at a time and this script does not walk — a real but
    narrow blind spot, distinct from and much smaller than "chart text is
    invisible to a file-reading auditor." Series/axis/legend/title-level
    fonts, which is where the vast majority of chart font defects live
    (per-point overrides are rare), are fully covered.
  - Does not check whether a declared font is installed on the machine
    running the audit — that is a rendering-environment property, not a
    property of the deck file, and is deliberately out of scope (a deck can
    target a different machine than the one auditing it).

Usage:
    python3 audit_deck.py path/to/deck.pptx
    python3 audit_deck.py path/to/deck.pptx --json   # machine-readable output
    python3 audit_deck.py path/to/deck.pptx --expect "Be Vietnam Pro,Manrope"
"""

import re
import sys
import json
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)

# ── Config: thresholds from deck-craft-rules.md (defaults, not hard gates) ──
TITLE_MAX_WORDS = 15
TITLE_MAX_LINES = 2
TITLE_MIN_FONT_PT = 28          # under 28pt on a title = text-stuffing red flag
TITLE_TOP_ZONE_RATIO = 0.30     # title candidates live in the top 30% of the slide
TABLE_MAX_ROWS_FOR_TABLE = 20   # tables preferred at <=20 numbers; beyond that, use a chart

# ── Font inventory: system-font hard flag — MECHANICAL, HIGH confidence ─────
# Office/OS default fonts. Their presence anywhere in a shipped institutional
# deliverable is unconditionally wrong — this is not a brand preference, it's
# a universal "don't ship the default" rule, independent of which brand's
# deck this is. Matched case-insensitively against python-pptx's font.name.
# Do NOT add brand-specific "wrong font" names here — that belongs in the
# optional --expect list, supplied by the caller from whichever brand skill
# is co-loaded. This set stays fixed and brand-agnostic.
SYSTEM_FONTS = {
    "calibri", "arial", "times new roman", "helvetica", "cambria", "verdana",
    "tahoma",
}

# Curated Vietnamese + English verb/copula indicators. Precision favored over
# recall on purpose — see "làm" exclusion note below. Extend this list as
# new decks surface false negatives; do not delete entries to fix one deck.
VERB_INDICATORS = [
    # Vietnamese copulas / core verbs
    "là", "có", "cần", "phải", "nên", "sẽ", "đang", "được", "bị",
    # Vietnamese action verbs seen across deck conventions
    "xác nhận", "cấm", "tăng", "giảm", "đọc", "phục vụ", "chốt", "thay đổi",
    "đề nghị", "phê duyệt", "yêu cầu", "mở rộng", "dẫn", "xây dựng",
    "kiểm soát", "đạt", "vượt", "chặn", "khoá", "cắt", "chuyển", "chi",
    "mua", "bán", "ký", "triển khai", "khởi động", "hoàn thành", "báo cáo",
    "theo dõi", "tối ưu", "phân bổ", "ưu tiên", "quyết định", "chọn",
    # NOTE: "làm" deliberately excluded — it is a function-verb fragment
    # embedded in many nominal compounds ("làm chủ", "công việc làm") and
    # produces false-positive verb hits on topic-label titles. If a title's
    # only verb-shaped token is "làm", it still counts as a candidate FAIL —
    # confirmed against the "Từ hợp tác toàn cầu đến làm chủ công nghệ"
    # (topic-label, correctly fails) vs "Từ 45 triệu VNĐ ... đến ~602.000
    # tiếp cận thực tế" (data-anchored, correctly passes via the
    # quantitative-anchor exemption below, not via "làm").
    # English action/copula verbs (for EN or mixed-language decks)
    "is", "are", "was", "were", "will", "must", "needs", "requires",
    "confirms", "increases", "decreases", "reduces", "grows", "cuts",
    "requests", "recommends", "approves", "targets", "drives", "wins",
    "loses", "delivers", "builds", "launches", "blocks", "unlocks",
]

# Quantitative-anchor exemption: a data-led headline ("50M VND ads in a
# 175M budget") functions as an implicit assertion even with no explicit
# verb. Require a UNIT-bearing number (currency/percent/grouped-thousands),
# not a bare digit — "Top 3 risks" has a digit but no unit, and still fails.
QUANT_ANCHOR_RE = re.compile(
    r"(VNĐ|đ\b|triệu|tỷ|%|\$|USD|\d{1,3}(?:[.,]\d{3})+|\d+\s*/\s*\d+)",
    re.IGNORECASE,
)

# Unresolved-conditional opener: "If X changes" with no resolution clause
# is a dangling hypothesis, not a complete assertion — fails even if it
# contains a verb ("thay đổi"/"changes").
CONDITIONAL_OPENER_RE = re.compile(
    r"^(Nếu|Khi|Trong trường hợp|If|When|In case)\b", re.IGNORECASE
)
CONDITIONAL_RESOLUTION_RE = re.compile(
    r"\b(thì|sẽ|cần|phải|then|will|must)\b", re.IGNORECASE
)

FOOTER_PAGENUM_RE = re.compile(r"^\s*\d{1,3}\s*/\s*\d{1,3}\s*$")

# ── Rule 8: reader-only content — process-meta marker scan ──────────────────
# HEURISTIC ONLY. Curated English + Vietnamese producer/reviewer-facing
# markers. The discriminator that actually decides whether a hit is a
# violation is WHO THE LINE IS ADDRESSED TO (the deck's real audience vs.
# the producer/an internal reviewer) — this script has no way to determine
# that from a string match, so every hit below is reported as a REVIEW
# CANDIDATE, never a verdict. See deck-narrative/SKILL.md Rule 8 for the
# addressed-to test and two worked KEEP examples ("Cần xác nhận: đội
# marketing đang kiểm soát các trang hiện có" and "Audit 07/2026 (chưa xác
# định)") that trip a marker here and are still correct as-is — a naive
# auto-strip on any of these strings would gut the deck's most important
# slides. Extend this list as new decks surface markers; do not delete
# entries to silence one false positive — fix it at the read-through step
# instead.
PROCESS_META_MARKERS = [
    # English — producer/reviewer scaffolding
    "tbd", "t.b.d", "to be determined", "to be confirmed", "placeholder",
    "needs review", "under review", "review needed", "open question",
    "pending review", "pending approval", "reviewer note", "internal note",
    "draft version", "confirm with", "awaiting confirmation",
    "comment on content", "decision needed",
    # Vietnamese — producer/reviewer scaffolding
    "cần xác nhận", "chưa xác định", "chưa rõ", "tạm thời", "chờ duyệt",
    "chờ phản hồi", "ghi chú nội bộ", "bản nháp", "góp ý", "cần xem lại",
    "xác nhận với",
]


def emu_to_in(v):
    return v / 914400 if v is not None else None


def _runs_from_text_frame(text_frame, loc_prefix, entries):
    """Append {"font": name_or_None, "location": str} for every run in a
    text frame. font is None when the run has no explicit font — it
    inherits from the theme (see UNSET_FONT handling)."""
    for p_idx, para in enumerate(text_frame.paragraphs):
        for r_idx, run in enumerate(para.runs):
            entries.append({"font": run.font.name, "location": f"{loc_prefix}.p[{p_idx}].r[{r_idx}]"})


def collect_font_entries(slide, shape_idx_base=""):
    """Font inventory. Walks every text RUN on the slide — shape text
    frames and table cells — plus every reachable font-bearing element in
    embedded charts: chart title runs, category/value axis tick-label
    fonts, legend font, and plot/series data-label fonts. One entry per
    run/element: {"font": name_or_None, "location": str}.

    Known blind spot (see module docstring): an INDIVIDUAL chart data
    point's label font override (point.data_label.font) is not walked —
    python-pptx exposes those one point at a time and this function does
    not iterate points. Series/axis/legend/title-level fonts — where the
    large majority of chart font defects actually live — ARE covered."""
    entries = []
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            _runs_from_text_frame(shape.text_frame, f"shape[{shape_idx}]", entries)
        if getattr(shape, "has_table", False):
            for r_idx, row in enumerate(shape.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    _runs_from_text_frame(cell.text_frame, f"shape[{shape_idx}].table[{r_idx}][{c_idx}]", entries)
        if getattr(shape, "has_chart", False):
            chart = shape.chart
            chart_loc = f"shape[{shape_idx}].chart"
            try:
                if chart.has_title and chart.chart_title.has_text_frame:
                    _runs_from_text_frame(chart.chart_title.text_frame, f"{chart_loc}.title", entries)
            except Exception:
                pass
            try:
                entries.append({"font": chart.category_axis.tick_labels.font.name,
                                 "location": f"{chart_loc}.category_axis.tick_labels"})
            except Exception:
                pass
            try:
                entries.append({"font": chart.value_axis.tick_labels.font.name,
                                 "location": f"{chart_loc}.value_axis.tick_labels"})
            except Exception:
                pass
            try:
                if chart.has_legend:
                    entries.append({"font": chart.legend.font.name, "location": f"{chart_loc}.legend"})
            except Exception:
                pass
            for plot_idx, plot in enumerate(chart.plots):
                try:
                    if plot.has_data_labels:
                        entries.append({"font": plot.data_labels.font.name,
                                         "location": f"{chart_loc}.plot[{plot_idx}].data_labels"})
                except Exception:
                    pass
                for series_idx, series in enumerate(plot.series):
                    try:
                        entries.append({"font": series.data_labels.font.name,
                                         "location": f"{chart_loc}.plot[{plot_idx}].series[{series_idx}].data_labels"})
                    except Exception:
                        pass
    return entries


def max_font_pt(shape):
    best = 0.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                best = max(best, run.font.size.pt)
    return best


def find_title_shape(slide, slide_height_emu):
    """Largest-font text box in the top zone of the slide, excluding
    page-number footers. Matches the kicker+headline pattern shared by
    common deck-building skills (e.g. skill-tekki-strategic-deck)."""
    zone_limit = slide_height_emu * TITLE_TOP_ZONE_RATIO
    best_shape, best_size = None, 0.0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or FOOTER_PAGENUM_RE.match(text):
            continue
        if shape.top is None or shape.top >= zone_limit:
            continue
        size = max_font_pt(shape)
        if size > best_size:
            best_size, best_shape = size, shape
    return best_shape, best_size


def collect_text_blocks(slide):
    """Yield (text, location) for every text-bearing element on the slide —
    shape text frames AND table cells. Rule 8's process-meta scan needs body
    text and table cells, not just the title (a data-gap flag like "chưa xác
    định" typically lives in a table cell, not a headline)."""
    blocks = []
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                blocks.append((text, f"shape[{shape_idx}]"))
        if getattr(shape, "has_table", False):
            for r_idx, row in enumerate(shape.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    cell_text = cell.text_frame.text.strip()
                    if cell_text:
                        blocks.append((cell_text, f"shape[{shape_idx}].table[{r_idx}][{c_idx}]"))
    return blocks


def process_meta_scan(slide):
    """Rule 8 (reader-only content) — HEURISTIC, REVIEW CANDIDATES ONLY.
    Scans every text block on the slide for PROCESS_META_MARKERS. Returns a
    list of {"location", "marker", "text"} dicts. This function does NOT and
    CANNOT determine whether a hit is a violation — that requires the
    addressed-to test (who the line speaks to: the deck's audience, or the
    producer/an internal reviewer), which needs human/agent reading of
    context. A hit here is a candidate for review, not a fail. See module
    docstring and deck-narrative/SKILL.md Rule 8 before treating any entry
    in the returned list as something to remove."""
    hits = []
    for text, location in collect_text_blocks(slide):
        lowered = text.lower()
        for marker in PROCESS_META_MARKERS:
            if marker in lowered:
                hits.append({"location": location, "marker": marker, "text": text})
    return hits


def verb_test(title):
    """Returns (passed: bool, reasons: list[str]). Heuristic — see module
    docstring. Reasons are review notes, not proof."""
    reasons = []

    if CONDITIONAL_OPENER_RE.match(title.strip()) and not CONDITIONAL_RESOLUTION_RE.search(title):
        reasons.append("unresolved conditional opener — dependent clause only, no resolution")

    has_verb = any(
        re.search(rf"(?<!\w){re.escape(v)}(?!\w)", title, re.IGNORECASE)
        for v in VERB_INDICATORS
    )
    has_anchor = bool(QUANT_ANCHOR_RE.search(title))

    if not has_verb and not has_anchor:
        reasons.append("no verb detected and no quantitative anchor — reads as a topic label")

    return (len(reasons) == 0), reasons


def estimate_wrapped_lines(shape, text):
    """Rough estimate only — no rendering engine available. Uses shape
    width + font size to guess chars-per-line. Flag as ESTIMATED in output;
    a human/agent visual check (or a screenshot) is the real verifier."""
    if shape is None or shape.width is None:
        return None
    width_pt = emu_to_in(shape.width) * 72
    size_pt = max_font_pt(shape) or 24
    # ~0.55 average char-width-to-font-size ratio for proportional sans/serif
    chars_per_line = max(1, int(width_pt / (size_pt * 0.55)))
    explicit_lines = text.count("\n") + 1
    longest_para = max((len(p) for p in text.split("\n")), default=0)
    wrapped_for_longest = -(-longest_para // chars_per_line)  # ceil
    return max(explicit_lines, wrapped_for_longest)


def audit_slide(slide, idx, slide_height_emu, is_title_slide):
    title_shape, title_font = find_title_shape(slide, slide_height_emu)
    title_text = title_shape.text_frame.text.strip() if title_shape else ""
    title_text_flat = " ".join(title_text.split())

    word_count = len(title_text_flat.split()) if title_text_flat else 0
    explicit_lines = title_text.count("\n") + 1 if title_text else 0
    est_lines = estimate_wrapped_lines(title_shape, title_text) if title_text else None

    slide_max_font = 0.0
    table_count = 0
    chart_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            slide_max_font = max(slide_max_font, max_font_pt(shape))
        if getattr(shape, "has_table", False):
            table_count += 1
        if getattr(shape, "has_chart", False):
            chart_count += 1

    verb_ok, verb_reasons = (True, []) if is_title_slide or not title_text else verb_test(title_text_flat)
    process_meta_candidates = process_meta_scan(slide)
    font_entries = collect_font_entries(slide)

    return {
        "slide": idx,
        "title": title_text_flat,
        "title_font_pt": title_font,
        "slide_max_font_pt": slide_max_font,
        "word_count": word_count,
        "explicit_line_breaks": explicit_lines,
        "estimated_wrapped_lines": est_lines,
        "words_over_limit": word_count > TITLE_MAX_WORDS,
        "font_under_28pt": bool(title_text) and title_font < TITLE_MIN_FONT_PT,
        "verb_test_pass": verb_ok,
        "verb_test_reasons": verb_reasons,
        "table_count": table_count,
        "chart_count": chart_count,
        "is_title_slide": is_title_slide,
        "process_meta_candidates": process_meta_candidates,
        "font_entries": font_entries,
    }


def audit_deck(path, expect_families=None):
    """expect_families: optional list of brand font family names (caller-
    supplied, from whichever brand skill is co-loaded). When given, any
    non-system font not in this set is flagged as outside the expected
    family list. When omitted, the font inventory is reported without
    family-membership judgment beyond the always-on system-font check —
    this script stays brand-agnostic by default (see module docstring)."""
    prs = Presentation(path)
    slide_height = prs.slide_height
    results = []
    for i, slide in enumerate(prs.slides, 1):
        results.append(audit_slide(slide, i, slide_height, is_title_slide=(i == 1)))

    total_tables = sum(r["table_count"] for r in results)
    total_charts = sum(r["chart_count"] for r in results)
    verb_fails = [r["slide"] for r in results if not r["verb_test_pass"]]
    font_fails = [r["slide"] for r in results if r["font_under_28pt"]]
    word_fails = [r["slide"] for r in results if r["words_over_limit"]]
    process_meta_slides = [r["slide"] for r in results if r["process_meta_candidates"]]

    # ── Font inventory aggregation (deck-wide) ──
    font_inventory = {}          # family name -> run count (named fonts only)
    unset_font_count = 0
    system_font_hits = []        # [{slide, location, font}]
    expected_set_lower = {f.strip().lower() for f in (expect_families or [])}
    unexpected_family_hits = []  # [{slide, location, font}] — only populated if expect_families given

    for r in results:
        for entry in r["font_entries"]:
            fname = entry["font"]
            if fname is None:
                unset_font_count += 1
                continue
            font_inventory[fname] = font_inventory.get(fname, 0) + 1
            if fname.strip().lower() in SYSTEM_FONTS:
                system_font_hits.append({"slide": r["slide"], "location": entry["location"], "font": fname})
            elif expect_families and fname.strip().lower() not in expected_set_lower:
                unexpected_family_hits.append({"slide": r["slide"], "location": entry["location"], "font": fname})

    system_font_families = sorted({h["font"] for h in system_font_hits})
    slides_with_system_font = sorted({h["slide"] for h in system_font_hits})

    # Which of the --expect families actually showed up in the deck (found in
    # font_inventory, matched case-insensitively). Reported separately from
    # unexpected_family_hits so "0 of N expected families present" is visible
    # even when there happen to be zero unexpected-but-non-system fonts too
    # (e.g. a deck that is 100% system font has nothing left to flag as
    # "unexpected" under the elif branch above, which would otherwise print
    # as a false-clean "all named fonts are within --expect").
    expected_families_found = []
    if expect_families:
        inventory_lower = {f.strip().lower() for f in font_inventory}
        expected_families_found = sorted(
            f for f in expect_families if f.strip().lower() in inventory_lower
        )

    return {
        "path": str(path),
        "slide_count": len(results),
        "total_tables": total_tables,
        "total_charts": total_charts,
        "titles_failing_verb_test": verb_fails,
        "titles_under_28pt": font_fails,
        "titles_over_word_limit": word_fails,
        "slides_with_process_meta_candidates": process_meta_slides,
        "font_inventory": font_inventory,
        "unset_font_count": unset_font_count,
        "system_font_families": system_font_families,
        "slides_with_system_font": slides_with_system_font,
        "system_font_hits": system_font_hits,
        "expect_families": sorted(expect_families) if expect_families else None,
        "expected_families_found": expected_families_found,
        "unexpected_family_hits": unexpected_family_hits,
        "slides": results,
    }


def print_report(report):
    print(f"deck-narrative audit — {report['path']}")
    print(f"slides: {report['slide_count']}")
    print(f"tables: {report['total_tables']}   charts: {report['total_charts']}")
    print()
    print("PER-SLIDE MECHANICAL CHECKS")
    print("-" * 100)
    header = f"{'#':>3}  {'font':>5}  {'words':>5}  {'lines(explicit/est)':>20}  {'verb':>5}  {'tbl':>3}  {'chart':>5}  title"
    print(header)
    for r in report["slides"]:
        verb_flag = "  --" if r["is_title_slide"] else ("PASS" if r["verb_test_pass"] else "FAIL")
        lines_str = f"{r['explicit_line_breaks']}/{r['estimated_wrapped_lines']}"
        title_disp = r["title"][:70] + ("…" if len(r["title"]) > 70 else "")
        print(f"{r['slide']:>3}  {r['title_font_pt']:>5.1f}  {r['word_count']:>5}  {lines_str:>20}  "
              f"{verb_flag:>5}  {r['table_count']:>3}  {r['chart_count']:>5}  {title_disp}")

    print()
    print("FONT INVENTORY (deck-wide, mechanical)")
    print("-" * 100)
    total_named_runs = sum(report["font_inventory"].values())
    if total_named_runs == 0 and report["unset_font_count"] == 0:
        print("  (no text runs found)")
    else:
        for fname, count in sorted(report["font_inventory"].items(), key=lambda kv: -kv[1]):
            flag = "  <-- SYSTEM FONT" if fname.strip().lower() in SYSTEM_FONTS else ""
            print(f"  {fname}: {count} runs/elements{flag}")
        if report["unset_font_count"]:
            print(f"  (unset — theme-inherited): {report['unset_font_count']} runs/elements")

    print()
    print("FLAGS (mechanical, HIGH confidence)")
    print(f"  Titles under 28pt (text-stuffing red flag): {report['titles_under_28pt'] or 'none'}")
    print(f"  Titles over {TITLE_MAX_WORDS}-word limit: {report['titles_over_word_limit'] or 'none'}")
    if report["system_font_families"]:
        print(f"  SYSTEM FONT DETECTED — always wrong for a shipped institutional deck:")
        print(f"    families: {report['system_font_families']}")
        print(f"    slides affected: {report['slides_with_system_font']}")
        print(f"    total hits: {len(report['system_font_hits'])}")
    else:
        print("  System fonts (Calibri/Arial/Times New Roman/Helvetica/Cambria/Verdana/Tahoma): none found")
    if report["expect_families"] is not None:
        found = report["expected_families_found"]
        missing = [f for f in report["expect_families"] if f not in found]
        print(f"  --expect {report['expect_families']}: "
              f"{len(found)}/{len(report['expect_families'])} expected families found in deck")
        if missing:
            print(f"    NOT found in deck: {missing}")
        if report["unexpected_family_hits"]:
            unexpected_families = sorted({h["font"] for h in report["unexpected_family_hits"]})
            unexpected_slides = sorted({h["slide"] for h in report["unexpected_family_hits"]})
            print(f"    OUTSIDE expected set (non-system, non-expected):")
            print(f"      families: {unexpected_families}")
            print(f"      slides affected: {unexpected_slides}")
        elif not found:
            print(f"    No expected-family text found at all — every named font in this deck "
                  f"is either a system font (see above) or otherwise outside --expect.")
        else:
            print(f"    No fonts outside the expected set (beyond any system-font hits above)")
    print()
    print("FLAGS (heuristic — review candidates, not verdicts)")
    print(f"  Titles failing verb test: {report['titles_failing_verb_test'] or 'none'}")
    for r in report["slides"]:
        if r["verb_test_reasons"]:
            print(f"    slide {r['slide']}: {'; '.join(r['verb_test_reasons'])}")

    print()
    print("  Rule 8 TEST 1 — process-meta marker hits (REVIEW CANDIDATES ONLY —")
    print("  NOT violations; apply the addressed-to test before acting on any of these):")
    print(f"    slides with a hit: {report['slides_with_process_meta_candidates'] or 'none'}")
    for r in report["slides"]:
        for hit in r["process_meta_candidates"]:
            text_disp = hit["text"][:80] + ("…" if len(hit["text"]) > 80 else "")
            print(f"    slide {r['slide']} [{hit['location']}] marker='{hit['marker']}': {text_disp}")

    print()
    print("TITLE SEQUENCE — for horizontal-logic read-through (human/agent judgment required)")
    print("-" * 100)
    for r in report["slides"]:
        print(f"  {r['slide']:>3}. {r['title']}")

    print()
    print("NOT CHECKED BY THIS SCRIPT (judgment — see deck-narrative/SKILL.md):")
    print("  - 'so what' test (does the title state an implication, not just describe data)")
    print("  - one-message-per-slide (5s/10s tests, does-this-prove-the-title deletion rule)")
    print("  - Minto sequencing (conclusion on slide 2, slides 1-3 = full recommendation)")
    print("  - horizontal logic (titles read as a persuasive essay in sequence)")
    print("  - vertical logic (body proves title)")
    print("  - decision-slide framework fit (Challenge>Options>Decision, etc.)")
    print("  - pre-wire verification, footnote citation completeness")
    print("  - Rule 8 TEST 1 verdict (is a process-meta hit actually producer-facing —")
    print("    script only finds the marker string, never the addressee)")
    print("  - Rule 8 TEST 2, see-vs-say (does the audience need to SEE this content, or")
    print("    can the presenter SAY it live — pure judgment, applies even to slides that")
    print("    pass every other check; e.g. three single-decision slides that should")
    print("    compress to one, with supporting research spoken instead of shown)")
    print("  - per-data-point chart label font overrides (point.data_label.font) — a real")
    print("    but narrow blind spot; series/axis/legend/title-level chart fonts ARE covered")
    print("  - whether a declared font is installed on the auditing machine (rendering-")
    print("    environment property, not a property of the deck file — out of scope)")


def main():
    ap = argparse.ArgumentParser(description="Mechanical audit for deck-narrative rules.")
    ap.add_argument("pptx_path", type=str, help="Path to a .pptx deck")
    ap.add_argument("--json", action="store_true", help="Output JSON instead of a text report")
    ap.add_argument("--expect", type=str, default=None,
                     help="Comma-separated expected brand font families (e.g. "
                          "'Be Vietnam Pro,Manrope'). Optional — brand-specific, "
                          "supplied by the caller from whichever brand skill is "
                          "co-loaded. Without it, only the fixed system-font check runs.")
    args = ap.parse_args()

    path = Path(args.pptx_path).expanduser()
    if not path.exists():
        print(f"ERROR: file not found: {path}", file=sys.stderr)
        sys.exit(1)

    expect_families = [f.strip() for f in args.expect.split(",") if f.strip()] if args.expect else None

    report = audit_deck(path, expect_families=expect_families)
    if args.json:
        print(json.dumps(report, indent=2, ensure_ascii=False))
    else:
        print_report(report)


if __name__ == "__main__":
    main()
